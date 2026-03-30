"""PPTX presentation parser using markitdown and python-pptx."""

import io
import logging
from typing import Optional

from agent_rag.ingestion.parsing.base import BaseParser, ParsedDocument, ParsedImage
from agent_rag.ingestion.parsing.utils import (
    extract_title_from_text,
    normalize_text,
)

logger = logging.getLogger(__name__)


class PPTXParser(BaseParser):
    """
    Parser for Microsoft PowerPoint .pptx files.

    Uses markitdown library for robust PPTX to markdown conversion,
    with python-pptx fallback for advanced extraction.

    Supports:
    - .pptx, .ppt files
    - application/vnd.openxmlformats-officedocument.presentationml.presentation

    Priority: 10 (default for office documents)
    """

    def supports(self, source_type: str, extension: str, mime_type: str) -> bool:
        """Support .pptx and .ppt files."""
        if extension.lower() in ["pptx", "ppt"]:
            return True
        if mime_type and ("presentationml" in mime_type or "powerpoint" in mime_type):
            return True
        return False

    def parse(self, content: bytes, filename: str) -> ParsedDocument:
        """
        Parse PPTX file using markitdown.

        Args:
            content: Raw PPTX bytes
            filename: Original filename

        Returns:
            ParsedDocument with extracted text and metadata

        Raises:
            ImportError: If markitdown not installed
            Exception: If parsing fails
        """
        try:
            from markitdown import MarkItDown
        except ImportError:
            raise RuntimeError(
                "markitdown not installed. Run: pip install markitdown"
            )

        # Create temporary file-like object
        file_stream = io.BytesIO(content)

        # Convert to markdown
        md = MarkItDown()

        try:
            result = md.convert_stream(file_stream, file_extension=".pptx")
        except Exception as e:
            logger.error(f"markitdown conversion failed: {e}")
            # Fallback to basic extraction if available
            return self._fallback_parse(content, filename)

        # Extract text
        text = result.text_content if hasattr(result, "text_content") else str(result)

        # Normalize text
        text = normalize_text(text)

        # Extract metadata
        metadata = {
            "filename": filename,
            "source_type": "file",
            "format": "pptx",
        }

        # Try to extract title
        title = extract_title_from_text(text)
        if title:
            metadata["title"] = title

        # markitdown may provide additional metadata
        if hasattr(result, "metadata") and result.metadata:
            metadata.update(result.metadata)

        # Extract images (markitdown may include image references)
        images = self._extract_images(content, filename)

        # Count slides
        slide_count = self._count_slides(content)
        if slide_count:
            metadata["slide_count"] = slide_count

        return ParsedDocument(
            text=text,
            metadata=metadata,
            images=images,
            links=[],  # Could extract from markdown links
            tables=[],  # markitdown converts tables to markdown
        )

    def _fallback_parse(self, content: bytes, filename: str) -> ParsedDocument:
        """
        Fallback parser using python-pptx.

        Args:
            content: Raw PPTX bytes
            filename: Original filename

        Returns:
            ParsedDocument with basic text extraction
        """
        try:
            from pptx import Presentation
        except ImportError:
            raise RuntimeError(
                "python-pptx not installed. Run: pip install python-pptx"
            )

        try:
            # Load presentation
            file_stream = io.BytesIO(content)
            prs = Presentation(file_stream)

            # Extract text from slides
            all_text = []

            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = []
                slide_text.append(f"## Slide {slide_num}")

                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())

                    # Handle tables
                    if shape.has_table:
                        table = shape.table
                        for row in table.rows:
                            row_text = " | ".join(
                                cell.text.strip() for cell in row.cells
                            )
                            if row_text.strip():
                                slide_text.append(row_text)

                if len(slide_text) > 1:  # More than just the header
                    all_text.append("\n".join(slide_text))

            text = "\n\n".join(all_text)
            text = normalize_text(text)

            # Basic metadata
            metadata = {
                "filename": filename,
                "source_type": "file",
                "format": "pptx",
                "parser": "fallback_python_pptx",
                "slide_count": len(prs.slides),
            }

            # Extract core properties if available
            if hasattr(prs, "core_properties"):
                props = prs.core_properties
                if props.title:
                    metadata["title"] = props.title
                if props.author:
                    metadata["author"] = props.author
                if props.created:
                    metadata["created"] = props.created.isoformat()
                if props.modified:
                    metadata["modified"] = props.modified.isoformat()

            return ParsedDocument(
                text=text,
                metadata=metadata,
                images=[],
                links=[],
                tables=[],
            )

        except Exception as e:
            logger.error(f"Fallback PPTX parsing failed: {e}")
            raise RuntimeError(f"Failed to parse PPTX file: {e}")

    def _extract_images(
        self,
        content: bytes,
        filename: str,
    ) -> list[ParsedImage]:
        """
        Extract images from PPTX file.

        Args:
            content: Raw PPTX bytes
            filename: Original filename

        Returns:
            List of ParsedImage objects
        """
        images = []

        try:
            from pptx import Presentation
            from pptx.enum.shapes import MSO_SHAPE_TYPE
        except ImportError:
            return images

        try:
            file_stream = io.BytesIO(content)
            prs = Presentation(file_stream)

            image_count = 0
            for slide_num, slide in enumerate(prs.slides, 1):
                for shape in slide.shapes:
                    # Check for picture shapes
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        try:
                            image = shape.image
                            image_bytes = image.blob
                            content_type = image.content_type or "image/png"

                            # Determine extension
                            ext = content_type.split("/")[-1]
                            if ext == "jpeg":
                                ext = "jpg"

                            image_count += 1
                            image_id = f"slide{slide_num}_img{image_count}"

                            images.append(
                                ParsedImage(
                                    image_id=image_id,
                                    content=image_bytes,
                                    mime_type=content_type,
                                    page_number=slide_num,
                                    caption=None,
                                )
                            )

                        except Exception as e:
                            logger.warning(
                                f"Failed to extract image from slide {slide_num}: {e}"
                            )
                            continue

        except Exception as e:
            logger.warning(f"Image extraction from PPTX failed: {e}")

        return images

    def _count_slides(self, content: bytes) -> Optional[int]:
        """
        Count number of slides in presentation.

        Args:
            content: Raw PPTX bytes

        Returns:
            Number of slides or None if cannot determine
        """
        try:
            from pptx import Presentation

            file_stream = io.BytesIO(content)
            prs = Presentation(file_stream)
            return len(prs.slides)
        except Exception:
            return None

    @property
    def priority(self) -> int:
        """Default priority for office documents."""
        return 10
